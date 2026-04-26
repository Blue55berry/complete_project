import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
from pptx import Presentation

prs = Presentation(r'c:\dev\flutter_pro\RiskGaurd1\project final\RiskGuard Final Review.pptx')
print(f"Total slides: {len(prs.slides)}")
for i, slide in enumerate(prs.slides):
    title = ""
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text:
            title = shape.text.replace('\n', ' | ')[:80]
            break
    print(f"  Slide {i+1}: {title}")
