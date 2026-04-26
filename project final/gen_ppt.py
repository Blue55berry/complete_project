"""Generate RiskGuard Final Review PPT - copy slide1 from review-1, rest follows Zenthor format."""
import copy, io, os
from lxml import etree
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

BASE = r'c:\dev\flutter_pro\RiskGaurd1\project final'

# Step 1: Copy Zenthor as base template
import shutil
OUT = os.path.join(BASE, 'RiskGuard - Final Review.pptx')
shutil.copy2(os.path.join(BASE, 'Zenthor Final Review.pptx'), OUT)

prs = Presentation(OUT)
rev1 = Presentation(os.path.join(BASE, 'RiskGaurd - review-1.pptx'))

# Step 2: Delete ALL Zenthor slides
xml_slides = prs.slides._sldIdLst
while len(xml_slides) > 0:
    rId = xml_slides[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs.part.drop_rel(rId)
    xml_slides.remove(xml_slides[0])

# Step 3: Copy slide 1 from review-1 as-is
src_slide = rev1.slides[0]
layout = prs.slide_layouts[0]
new_slide = prs.slides.add_slide(layout)

# Clear default shapes
for sh in list(new_slide.shapes):
    sh._element.getparent().remove(sh._element)

# Copy shapes from source slide
ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
      'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'}

for child in src_slide.shapes._spTree:
    tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
    if tag in ('sp', 'pic', 'grpSp', 'graphicFrame', 'cxnSp'):
        el = copy.deepcopy(child)
        # Fix image relationships
        for blip in el.findall('.//a:blip', ns):
            old_rId = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
            if old_rId and old_rId in src_slide.part.rels:
                img_part = src_slide.part.rels[old_rId].target_part
                new_rId = new_slide.part.relate_to(img_part, RT.IMAGE)
                blip.set('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed', new_rId)
        new_slide.shapes._spTree.append(el)

# Step 4: Helper to add content slides
CL = prs.slide_layouts[1]  # Title and Content

def add_slide(title, bullets):
    s = prs.slides.add_slide(CL)
    tf = s.placeholders[0].text_frame
    tf.clear()
    r = tf.paragraphs[0].add_run()
    r.text = title
    r.font.name = 'Times New Roman'
    body = s.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        r = p.add_run()
        r.text = b
        r.font.name = 'Times New Roman'
    return s

def add_empty_slide(title):
    """Slide with just a title - for screenshots/diagrams placeholder."""
    s = prs.slides.add_slide(CL)
    tf = s.placeholders[0].text_frame
    tf.clear()
    r = tf.paragraphs[0].add_run()
    r.text = title
    r.font.name = 'Times New Roman'
    # Clear content placeholder
    body = s.placeholders[1].text_frame
    body.clear()
    p = body.paragraphs[0]
    r = p.add_run()
    r.text = '[Screenshot/Diagram Placeholder]'
    r.font.name = 'Times New Roman'
    r.font.italic = True
    return s

def add_table_slide(title, header, rows):
    s = prs.slides.add_slide(CL)
    tf = s.placeholders[0].text_frame
    tf.clear()
    r = tf.paragraphs[0].add_run()
    r.text = title
    r.font.name = 'Times New Roman'
    sp = s.placeholders[1]._element
    sp.getparent().remove(sp)
    data = [header] + rows
    nr, nc = len(data), len(data[0])
    tbl = s.shapes.add_table(nr, nc, Emu(457200), Emu(1800000), Emu(8229600), Emu(4200000)).table
    for ri in range(nr):
        for ci in range(nc):
            cell = tbl.cell(ri, ci)
            cell.text = data[ri][ci]
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.name = 'Times New Roman'
                    run.font.size = Pt(13)
                    if ri == 0:
                        run.font.bold = True
    return s

# ============ SLIDES ============

# 2. Introduction
add_slide('Introduction', [
    'AI-generated media poses escalating threats to digital trust.',
    'Deepfake images, cloned voices, and AI text are increasingly realistic.',
    'Existing tools require manual upload and support only single modality.',
    'Citizens lack real-time protection during everyday app usage.',
    'RiskGuard addresses this with a multi-modal detection platform.',
    'Combines local signal processing with cloud-based AI models.',
])

# 3. Abstract
add_slide('Abstract', [
    'Multi-modal AI deepfake detection across image, voice, text, and video.',
    'Hybrid CPU + GPU architecture with 6-signal ensemble detection.',
    'Blockchain evidence preservation using SHA-256, IPFS, and Merkle trees.',
    'Flutter mobile app with real-time protection overlay.',
    'FastAPI web dashboard for law enforcement investigators.',
    'Achieves 70-92% image, 65-85% voice, 70-88% text accuracy.',
])

# 4. Problem Statement
add_slide('Problem Statement', [
    'AI voice cloning scams and impersonation attacks.',
    'Deepfake videos spreading misinformation.',
    'AI-generated phishing messages and social engineering.',
    'Lack of unified tools for detecting multiple media types.',
    'Privacy risks due to cloud-only detection systems.',
    'No proactive protection during browsing social media apps.',
])

# 5. Objectives
add_slide('Objectives', [
    'Build multi-signal ensemble detection for four modalities.',
    'Achieve 60-92% accuracy on commodity hardware.',
    'Implement blockchain evidence ledger with Merkle batching.',
    'Reduce blockchain costs by up to 99.7%.',
    'Develop real-time overlay monitoring whitelisted apps.',
    'Create dual interface for citizens and investigators.',
])

# 6. Existing System
add_slide('Existing System', [
    'Tools like Microsoft Video Authenticator, Deepware Scanner.',
    'Require manual upload of suspicious content.',
    'Limited to single-modality analysis.',
    'Cloud-based GPU inference with processing delays.',
    'Blockchain systems operate independently of detection.',
    'No real-time automated protection.',
])

# 7. Limitations of Existing System
add_slide('Limitations of Existing System', [
    'Manual user-initiated analysis only.',
    'Cannot cross-correlate signals across media types.',
    'Slow processing: 10-60 seconds per item.',
    'Requires constant cloud connectivity.',
    'Detection and evidence preservation are fragmented.',
    'No offline detection capability.',
    'No proactive monitoring of messaging apps.',
])

# 8. Proposed System
add_slide('Proposed System', [
    'Comprehensive multi-modal AI detection platform.',
    'Proactive real-time overlay monitors apps automatically.',
    'Hybrid CPU + GPU local and cloud detection.',
    '6-signal image, 6-signal voice, 4-signal text ensemble.',
    'Integrated blockchain evidence with Merkle batching.',
    'Operates on commodity 8 GB RAM hardware.',
    'Dual interface: Flutter app + web dashboard.',
])

# 9. Advantages of Proposed System
add_slide('Advantages of Proposed System', [
    'Multi-modal detection in a single unified platform.',
    'Real-time proactive protection through system overlay.',
    'Hybrid architecture enables offline detection.',
    'Merkle batching reduces gas costs by 99.7%.',
    'Live call monitoring with voice deepfake detection.',
    'URL verification with threat intelligence feeds.',
    'Processing under 200ms for voice, under 150ms for images.',
])

# 10. System Architecture - Overview
add_slide('System Architecture - Overview', [
    'Three-tier: Flutter frontend, FastAPI backend, external services.',
    'Mobile app communicates via HTTPS REST APIs.',
    'Backend processes AI detection across four modalities.',
    'External: HuggingFace, Colab, Pinata IPFS, Polygon blockchain.',
    'Overlay runs independently via Android Accessibility Services.',
])

# 11. Architecture Diagram (placeholder)
add_empty_slide('System Architecture Diagram')

# 12. Image Analysis Module
add_slide('Image Analysis Module', [
    '6-signal type-adaptive weighted ensemble.',
    'Classifies images into photograph, digital art, or screenshot.',
    'NPR extracts camera sensor noise patterns.',
    'DCT spectral analysis detects GAN/diffusion fingerprints.',
    'Haar wavelet decomposition for texture features.',
    'Cloud CNN + HuggingFace classifiers.',
    'Failed signals excluded, weights redistributed.',
])

# 13. Voice Analysis Module
add_slide('Voice Analysis Module', [
    '6-signal hybrid CPU + GPU architecture.',
    'Real-time streaming: 0.5s chunks under 200ms.',
    'LFCC (30%) - Cepstral coefficient analysis.',
    'CQT Phase (20%) - Phase coherence detection.',
    'Modulation (20%) - Temporal envelope analysis.',
    'Pitch/F0 (20%) - Autocorrelation contour tracking.',
    'wav2vec2 (40%) - Fine-tuned on ASVspoof2019.',
])

# 14. Text Analysis Module
add_slide('Text Analysis Module', [
    '4-signal ensemble with cloud and local methods.',
    'DeBERTa (45%) - Fine-tuned on ChatGPT outputs.',
    'RoBERTa (15%) - Short-text dampening fix.',
    'Binoculars (20%) - Zero-shot LLM perplexity proxy.',
    'Local Statistics (20%) - Lexical diversity analysis.',
    'Independent phishing analysis with URL detection.',
    'Risk score 0-100 with explanation.',
])

# 15. Video Analysis Module
add_slide('Video Analysis Module', [
    '2-signal temporal-aware pipeline.',
    'Per-frame image signal (60%) - Full 6-signal pipeline.',
    'Temporal coherence (40%) - Farneback optical flow.',
    'Samples at 3 FPS, maximum 30 frames.',
    'Deepfakes show abrupt flow variance.',
    'Real video shows smooth optical flow.',
])

# 16. Blockchain Evidence Module
add_slide('Blockchain Evidence Module', [
    'Follows Immutable Chain of Custody principle.',
    'SHA-256 hash fingerprints the evidence file.',
    'IPFS upload via Pinata for decentralized storage.',
    'SQLite records off-chain metadata.',
    'Merkle tree groups hashes, only root stored on-chain.',
    'Single Polygon transaction anchors entire batch.',
    'Cost reduction: 99.7% for 100 records.',
])

# 17. Real-Time Protection Overlay
add_slide('Real-Time Protection Overlay', [
    '2300+ line Android system overlay.',
    'Accessibility Services monitors whitelisted apps.',
    'Supports WhatsApp, Chrome, Instagram, Telegram, Facebook.',
    'Three modes: bubble, verdict card, call monitoring chip.',
    'URL verification against URLhaus and blacklists.',
    'Media capture for deepfake analysis.',
    'Live call monitoring during active phone calls.',
])

# 18. Technologies Used
add_table_slide('Technologies Used',
    ['Component', 'Technology', 'Purpose'],
    [['Mobile Frontend', 'Flutter / Dart', 'Cross-platform app with overlay'],
     ['Backend API', 'FastAPI (Python 3.11+)', 'Async REST API'],
     ['AI Detection', 'NumPy, SciPy, OpenCV', 'Signal processing'],
     ['Cloud Models', 'HuggingFace API', 'Transformer classification'],
     ['GPU Inference', 'Google Colab + ONNX', 'wav2vec2, CNN classifier'],
     ['Blockchain', 'web3.py, Polygon Amoy', 'Evidence anchoring'],
     ['Smart Contract', 'Solidity', 'Merkle root storage'],
     ['IPFS', 'Pinata API', 'Decentralized storage'],
     ['Dashboard', 'FastAPI + Jinja2 + SSE', 'Investigator interface'],
     ['State Mgmt', 'Provider (Flutter)', 'Reactive state management']])

# 19. Working of the System
add_slide('Working of the System', [
    'Image: Classify type, assign weights, run 6 signals concurrently.',
    'Voice: Preprocess audio, run 6 signals, fuse with cloud model.',
    'Text: Run 4-signal ensemble with phishing analysis layer.',
    'Overlay: Detect app, activate bubble, scan content automatically.',
    'Evidence: Hash file, upload IPFS, Merkle batch, anchor on-chain.',
    'Dashboard: SSE live feed, toast alerts, batch anchoring.',
])

# 20. Testing and Evaluation
add_slide('Testing and Evaluation', [
    'Unit Testing - Individual modules validated.',
    'Integration Testing - Frontend-backend communication verified.',
    'System Testing - Full application with real-world inputs.',
    'User Acceptance Testing - End-user usability evaluation.',
    'Performance: response time, speed, reliability assessed.',
    'Feedback-driven UI enhancements.',
])

# 21. Test Cases
add_table_slide('Test Case Results',
    ['Test ID', 'Input Type', 'Expected Output', 'Result'],
    [['TC01', 'Real Image', 'Classified as Real', 'Passed'],
     ['TC02', 'Deepfake Image', 'Classified as Fake', 'Passed'],
     ['TC03', 'Real Video', 'Real with Confidence', 'Passed'],
     ['TC04', 'Deepfake Video', 'Fake with Confidence', 'Passed'],
     ['TC05', 'Invalid File', 'Error Message', 'Passed']])

# 22. Results and Performance
add_slide('Results and Performance', [
    'Image detection: 70-92% accuracy on photographs.',
    'Type-adaptive weighting prevents false positives on art.',
    'Voice detection: 65-85% accuracy on audio clips.',
    'Real-time streaming under 200ms per chunk.',
    'Text detection: 70-88% on long-form content.',
    'Local image analysis under 150ms.',
    'Merkle batching reduces costs by 99.7%.',
])

# 23-25. Screenshot placeholders
add_empty_slide('Mobile Application Screenshots')
add_empty_slide('Real-Time Protection Overlay Screenshots')
add_empty_slide('Web Dashboard & Backend Screenshots')

# 26. Future Scope
add_slide('Future Scope', [
    'Deploy quantized models locally via ONNX Runtime.',
    'Adversarial robustness testing against evasion attacks.',
    'Cross-modal fusion for multi-modal documents.',
    'Migrate from Polygon Amoy testnet to mainnet.',
    'Extend overlay system to iOS.',
    'On-device NLP for SMS phishing detection.',
    'Continuous retraining for new AI techniques.',
])

# 27. Conclusion
add_slide('Conclusion', [
    'RiskGuard integrates detection, evidence, and real-time protection.',
    'Production-grade detection on commodity hardware.',
    'Type-adaptive image ensemble prevents false positives.',
    'Voice pipeline achieves real-time under 200ms.',
    'Blockchain costs reduced 99.7% via Merkle batching.',
    'Overlay provides continuous background monitoring.',
    'Dual interface serves citizens and law enforcement.',
])

# 28. References
add_slide('References', [
    'Cozzolino & Verdoliva (2020) - Noiseprint. IEEE TIFS.',
    'Tak et al. (2021) - End-to-End Anti-Spoofing. ICASSP.',
    'Baevski et al. (2020) - wav2vec 2.0. NeurIPS.',
    'Hans et al. (2024) - Spotting LLMs with Binoculars. ICML.',
    'Sahidullah et al. (2015) - Block Level Features.',
    'ASVspoof Consortium (2024) - ASVspoof 5. Interspeech.',
    'Benet, J. (2014) - IPFS. arXiv:1407.3561.',
])

# 29. Conference Certificates (placeholder)
add_empty_slide('Conference Certificates')

# 30. Thank You
s = prs.slides.add_slide(CL)
tf = s.placeholders[0].text_frame
tf.clear()
sp = s.placeholders[1]._element
sp.getparent().remove(sp)
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = 'THANK YOU'
r.font.name = 'Times New Roman'
r.font.size = Pt(60)

prs.save(OUT)
print(f"Done! {len(prs.slides)} slides saved to {OUT}")
