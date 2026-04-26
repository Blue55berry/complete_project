import sys
import io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN

def get_color_info(color):
    try:
        if color and color.type:
            try:
                return str(color.rgb)
            except:
                return f"type={color.type}"
        return None
    except:
        return None

def analyze_pptx(filepath, label):
    prs = Presentation(filepath)
    print(f"\n{'='*80}")
    print(f"ANALYZING: {label}")
    print(f"{'='*80}")
    print(f"Slide width: {prs.slide_width}, height: {prs.slide_height}")
    print(f"Total slides: {len(prs.slides)}")
    
    for slide_idx, slide in enumerate(prs.slides):
        print(f"\n--- SLIDE {slide_idx + 1} ---")
        layout = slide.slide_layout
        print(f"  Layout: {layout.name}")
        
        for shape_idx, shape in enumerate(slide.shapes):
            print(f"\n  Shape {shape_idx}: type={shape.shape_type}, name='{shape.name}'")
            print(f"    Pos: left={shape.left}, top={shape.top}, w={shape.width}, h={shape.height}")
            
            if hasattr(shape, "text") and shape.text:
                print(f"    Text: '{shape.text[:300]}'")
            
            if shape.has_text_frame:
                for para_idx, para in enumerate(shape.text_frame.paragraphs):
                    align = para.alignment
                    for run in para.runs:
                        font = run.font
                        color_info = get_color_info(font.color)
                        size_pt = round(font.size / 12700, 1) if font.size else None
                        print(f"      Run: '{run.text[:80]}' | font={font.name} size={size_pt}pt bold={font.bold} italic={font.italic} color={color_info} align={align}")
            
            if shape.shape_type == 13:
                print(f"    [IMAGE] {shape.image.content_type}")

analyze_pptx(r'c:\dev\flutter_pro\RiskGaurd1\project final\Zenthor Final Review.pptx', 'Zenthor Final Review')
print("\n\n" + "X"*80 + "\n\n")
analyze_pptx(r'c:\dev\flutter_pro\RiskGaurd1\project final\RiskGaurd - review-1.pptx', 'RiskGuard Review 1')
