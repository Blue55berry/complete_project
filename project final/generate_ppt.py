"""Generate RiskGuard Final Review PPT using Zenthor as template."""
import copy, os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

BASE = r'c:\dev\flutter_pro\RiskGaurd1\project final'
ZENTHOR = os.path.join(BASE, 'Zenthor Final Review.pptx')
REVIEW1 = os.path.join(BASE, 'RiskGaurd - review-1.pptx')
OUTPUT = os.path.join(BASE, 'RiskGuard Final Review.pptx')

# Load Zenthor as template
prs = Presentation(ZENTHOR)

# Extract images from review-1 slide 1
rev1 = Presentation(REVIEW1)
rev1_images = {}
for shape in rev1.slides[0].shapes:
    if shape.shape_type == 13:  # Picture
        blob = shape.image.blob
        ct = shape.image.content_type
        rev1_images[shape.name] = {'blob': blob, 'ct': ct,
            'left': shape.left, 'top': shape.top,
            'width': shape.width, 'height': shape.height}

# Extract arch diagram from review-1 slide 8
arch_img = None
if len(rev1.slides) > 7:
    for shape in rev1.slides[7].shapes:
        if shape.shape_type == 13:
            arch_img = {'blob': shape.image.blob, 'ct': shape.image.content_type}

# Delete all Zenthor slides
while len(prs.slides._sldIdLst) > 0:
    rId = prs.slides._sldIdLst[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

# Get layouts
title_layout = prs.slide_layouts[0]  # Title Slide
content_layout = prs.slide_layouts[1]  # Title and Content

def add_title_content_slide(title, bullets, font_name='Times New Roman'):
    slide = prs.slides.add_slide(content_layout)
    # Set title
    tf = slide.placeholders[0].text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = font_name
    # Set bullets
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = body.paragraphs[0]
        else:
            p = body.add_paragraph()
        if isinstance(bullet, tuple):
            # (text, bold_flag)
            run = p.add_run()
            run.text = bullet[0]
            run.font.name = font_name
            run.font.bold = bullet[1]
        else:
            run = p.add_run()
            run.text = bullet
            run.font.name = font_name
    return slide

def add_image_slide(title, img_data, font_name='Times New Roman'):
    slide = prs.slides.add_slide(content_layout)
    tf = slide.placeholders[0].text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.name = font_name
    if img_data:
        import io
        slide.shapes.add_picture(
            io.BytesIO(img_data['blob']),
            Emu(283361), Emu(1295400), Emu(8577276), Emu(4682321))
    return slide

# ============================================================
# SLIDE 1: Title Slide (matching review-1 format in Zenthor theme)
# ============================================================
slide1 = prs.slides.add_slide(title_layout)

# College header rectangle
tf = slide1.placeholders[0].text_frame
tf.clear()
p = tf.paragraphs[0]
run = p.add_run()
run.text = 'RISKGUARD'
run.font.name = 'Times New Roman'

tf2 = slide1.placeholders[1].text_frame
tf2.clear()
p = tf2.paragraphs[0]
run = p.add_run()
run.text = 'A Multi-Modal AI Deepfake Detection System'
run.font.name = 'Times New Roman'

# College header box
from pptx.util import Emu
txBox = slide1.shapes.add_textbox(Emu(118242), Emu(93507), Emu(9025758), Emu(1200329))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = 'Erode Sengunthar Engineering College'
run.font.name = 'Times New Roman'
run.font.size = Pt(28)
p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
run2 = p2.add_run()
run2.text = '(An Autonomous Institution)'
run2.font.name = 'Times New Roman'
run2.font.size = Pt(24)
p3 = tf.add_paragraph()
p3.alignment = PP_ALIGN.CENTER
run3 = p3.add_run()
run3.text = 'Perundurai, Erode - 638057'
run3.font.name = 'Times New Roman'
run3.font.size = Pt(20)

# Add college logos from review-1
import io
for name, img in rev1_images.items():
    slide1.shapes.add_picture(
        io.BytesIO(img['blob']),
        img['left'], img['top'], img['width'], img['height'])

# Guide info
txBox2 = slide1.shapes.add_textbox(Emu(291662), Emu(4462273), Emu(4390066), Emu(1292662))
tf = txBox2.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
run = p.add_run()
run.text = 'UNDER THE GUIDANCE OF'
run.font.name = 'Times New Roman'
run.font.bold = True
run.font.size = Pt(18)
for line in ['Mr. R. Narendran, B.E., M.E., (Ph.D).,',
             'Assistant Professor,',
             'Department of Information Technology']:
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    run2 = p2.add_run()
    run2.text = line
    run2.font.name = 'Times New Roman'
    run2.font.size = Pt(20)

# Presented by
txBox3 = slide1.shapes.add_textbox(Emu(5038344), Emu(4462273), Emu(3657600), Emu(1920240))
tf = txBox3.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
run = p.add_run()
run.text = 'PRESENTED BY'
run.font.name = 'Times New Roman'
run.font.bold = True
run.font.size = Pt(18)
students = [
    'Abi Prasath R          730422205002',
    'Harija A                   730422205023',
    'Hemabala S P           730422205030',
]
for s in students:
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    run2 = p2.add_run()
    run2.text = s
    run2.font.name = 'Times New Roman'
    run2.font.size = Pt(20)

# ============================================================
# SLIDE 2: Introduction
# ============================================================
add_title_content_slide('Introduction', [
    'AI-generated media poses escalating threats to digital trust and public safety.',
    'Deepfake images, cloned voices, AI-authored phishing text are increasingly realistic.',
    'Traditional systems cannot distinguish authentic content from synthetic media.',
    'Law enforcement and citizens need proactive, real-time detection tools.',
    'RiskGuard provides a comprehensive multi-modal detection platform.',
    'Combines local signal processing with cloud-based AI models.',
])

# SLIDE 3: Abstract
add_title_content_slide('Abstract', [
    'RiskGuard is a multi-modal AI deepfake detection system.',
    'Detects synthetic content across image, voice, text, and video modalities.',
    'Uses hybrid CPU + GPU architecture with 6-signal ensemble detection.',
    'Blockchain evidence preservation via SHA-256, IPFS, and Merkle trees.',
    'Flutter mobile app with real-time protection overlay for citizens.',
    'FastAPI web dashboard for law enforcement investigators.',
    'Achieves 70-92% image, 65-85% voice, and 70-88% text detection accuracy.',
])

# SLIDE 4: Need for the System
add_title_content_slide('Need for the System', [
    'Rapid proliferation of AI-generated deepfake content across all media types.',
    'Existing tools require manual content submission for analysis.',
    'Single-modality detection leaves gaps in security coverage.',
    'No integrated system combining detection with evidence preservation.',
    'Citizens lack real-time protection during everyday app usage.',
    'Law enforcement needs tamper-proof forensic evidence chains.',
])

# SLIDE 5: Problem Statement
add_title_content_slide('Problem Statement', [
    'Accurately distinguishing AI-generated content across multiple modalities.',
    'Ensuring tamper-proof chain of custody for digital forensic evidence.',
    'Providing real-time detection without requiring manual content submission.',
    'Most existing solutions are limited to single modality (image OR video OR audio).',
    'Processing delays of 10-60 seconds per item in current tools.',
    'No proactive protection during browsing WhatsApp, Chrome, or Instagram.',
])

# SLIDE 6: Objectives
add_title_content_slide('Objectives', [
    'Build multi-signal ensemble detection for image, voice, text, and video.',
    'Achieve 60-92% accuracy without requiring GPU hardware on detection path.',
    'Implement blockchain evidence ledger with Merkle-batch anchoring.',
    'Reduce per-record blockchain costs by up to 99.7% via Merkle batching.',
    'Develop real-time protection overlay monitoring whitelisted applications.',
    'Create dual-interface: Flutter mobile app + FastAPI web dashboard.',
])

# SLIDE 7: Existing System
add_title_content_slide('Existing System', [
    'Microsoft Video Authenticator, Deepware Scanner, Sensity AI — standalone tools.',
    'Require users to manually upload suspicious content for verification.',
    'Limited to single-modality analysis (image OR video OR audio).',
    'Blockchain systems like Chainkit operate independently of detection.',
    'Cloud-based GPU inference with significant processing delays.',
    'No real-time, proactive protection during everyday app usage.',
])

# SLIDE 8: Limitations of Existing System
add_title_content_slide('Limitations of Existing System', [
    'Manual, user-initiated analysis — threats during normal usage go undetected.',
    'Single-modality — cannot cross-correlate signals across media types.',
    'Slow processing (10-60 seconds per item) requiring cloud connectivity.',
    'Blockchain evidence systems fragmented from detection workflow.',
    'No offline detection capability available.',
    'No proactive monitoring of messaging and social media applications.',
])

# SLIDE 9: Proposed System
add_title_content_slide('Proposed System', [
    'RiskGuard: comprehensive multi-modal AI detection platform.',
    'Proactive real-time overlay monitors whitelisted apps automatically.',
    'Hybrid CPU + GPU architecture for local + cloud detection.',
    '6-signal ensemble for images, 6-signal for voice, 4-signal for text.',
    'Integrated blockchain evidence chain with Merkle-batch anchoring.',
    'Operates on commodity 8 GB RAM hardware with optional cloud GPU.',
    'Dual interface serving citizens (Flutter) and investigators (web dashboard).',
])

# SLIDE 10: Advantages of Proposed System
add_title_content_slide('Advantages of Proposed System', [
    'Multi-modal detection across image, voice, text, and video in one platform.',
    'Real-time proactive protection through system overlay monitoring apps.',
    'Hybrid local + cloud architecture enabling offline detection.',
    'Merkle-batch anchoring reduces blockchain gas costs by 99.7%.',
    'Live call monitoring with voice deepfake detection during phone calls.',
    'URL verification with URLhaus threat feeds and community blacklists.',
    'Processing under 200ms for voice chunks, under 150ms for local images.',
])

# SLIDE 11: System Architecture - Overview
add_title_content_slide('System Architecture - Overview', [
    'Three-tier architecture: Flutter frontend, FastAPI backend, external services.',
    'Mobile app communicates via HTTPS REST APIs with backend.',
    'Backend processes AI detection across four modalities.',
    'External services: HuggingFace, Google Colab, Pinata IPFS, Polygon blockchain.',
    'System overlay runs independently via Android Accessibility Services.',
    'Platform channels bridge native Android services with Flutter UI.',
])

# SLIDE 12: Architecture Diagram
add_image_slide('System Architecture Diagram', arch_img)

# SLIDE 13: Image Analysis Module
add_title_content_slide('Image Analysis Module', [
    '6-signal type-adaptive weighted ensemble architecture.',
    'Classifies images into photograph, digital art, or screenshot.',
    'NPR (Noise Print Residual) detects camera sensor noise patterns.',
    'DCT spectral analysis detects GAN/diffusion frequency fingerprints.',
    'Haar wavelet decomposition analyzes multi-scale texture features.',
    'Cloud CNN + HuggingFace classifiers for deep learning signals.',
    'Graceful degradation: failed signals excluded, weights redistributed.',
])

# SLIDE 14: Image Detection Weights
slide14 = add_title_content_slide('Image Detection Signal Weights', [
    'Signal weights dynamically adjusted based on image type:',
])
# Add table
from pptx.util import Inches
table_data = [
    ['Signal', 'Photo Weight', 'Art Weight', 'Technique'],
    ['ONNX/Colab', '25%', '10%', 'CNN Binary Classifier'],
    ['Cloud HF', '25%', '10%', 'HuggingFace AI Detector'],
    ['NPR', '15%', '5%', 'Noise Print Residual'],
    ['Wavelet', '15%', '5%', 'Haar Wavelet Decomposition'],
    ['pHash', '10%', '5%', 'Perceptual Hashing (DCT)'],
    ['DCT', '10%', '65%', 'DCT Spectral Analysis'],
]
rows, cols = len(table_data), len(table_data[0])
tbl = slide14.shapes.add_table(rows, cols, Emu(457200), Emu(2400000), Emu(8229600), Emu(3200000)).table
for r in range(rows):
    for c in range(cols):
        cell = tbl.cell(r, c)
        cell.text = table_data[r][c]
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.name = 'Times New Roman'
                run.font.size = Pt(14)
                if r == 0:
                    run.font.bold = True

# SLIDE 15: Voice Analysis Module
add_title_content_slide('Voice Analysis Module', [
    '6-signal hybrid CPU + GPU architecture with real-time streaming.',
    'Processes 0.5-second chunks with analysis under 200ms per chunk.',
    'LFCC (30%) — Linear Frequency Cepstral Coefficients.',
    'CQT Phase Coherence (20%) — Constant-Q Transform analysis.',
    'Modulation Spectrum (20%) — Temporal envelope analysis.',
    'Pitch/F0 Contour (20%) — Autocorrelation-based tracking.',
    'Statistical Moments (10%) — Higher-order amplitude analysis.',
    'wav2vec2 GPU (40%) — Fine-tuned on ASVspoof2019 via ONNX.',
])

# SLIDE 16: Voice Detection Weights
slide16 = add_title_content_slide('Voice Detection Signal Architecture', [
    'Fusion: 60% local CPU + 40% Colab GPU (100% local when offline):',
])
vtbl_data = [
    ['Signal', 'Weight', 'Technique'],
    ['LFCC', '30%', 'Linear Frequency Cepstral Coefficients'],
    ['CQT Phase', '20%', 'Constant-Q Transform Phase Coherence'],
    ['Modulation', '20%', 'Temporal Envelope Modulation Spectrum'],
    ['Pitch/F0', '20%', 'Autocorrelation-based F0 Contour'],
    ['Statistical', '10%', 'Higher-order Amplitude Moments'],
    ['wav2vec2', '40%', 'wav2vec2-base ASVspoof2019 ONNX'],
]
rows, cols = len(vtbl_data), len(vtbl_data[0])
tbl = slide16.shapes.add_table(rows, cols, Emu(457200), Emu(2400000), Emu(8229600), Emu(3200000)).table
for r in range(rows):
    for c in range(cols):
        cell = tbl.cell(r, c)
        cell.text = vtbl_data[r][c]
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.name = 'Times New Roman'
                run.font.size = Pt(14)
                if r == 0:
                    run.font.bold = True

# SLIDE 17: Text Analysis Module
add_title_content_slide('Text Analysis Module', [
    '4-signal ensemble combining cloud transformers with local methods.',
    'DeBERTa (45%) — Fine-tuned on ChatGPT outputs for AI text detection.',
    'RoBERTa (15%) — Short-text dampening fix for casual messages.',
    'Binoculars Perplexity Proxy (20%) — Zero-shot LLM detection.',
    'Local Statistical Analysis (20%) — Lexical diversity and patterns.',
    'Independent phishing analysis layer with URL and urgency detection.',
    'Risk score (0-100) with explanation for each analyzed text.',
])

# SLIDE 18: Video Analysis Module
add_title_content_slide('Video Analysis Module', [
    '2-signal temporal-aware pipeline reusing full image engine per-frame.',
    'Per-frame image signal (60%) — Full 6-signal pipeline on sampled frames.',
    'Temporal coherence signal (40%) — Farneback dense optical flow.',
    'Samples at 3 FPS with maximum of 30 frames per video.',
    'Deepfakes show abrupt flow variance due to face swap discontinuities.',
    'Real video shows smooth, consistent optical flow patterns.',
])

# SLIDE 19: Blockchain Evidence Module
add_title_content_slide('Blockchain Evidence Module', [
    'Follows "Immutable Chain of Custody" principle.',
    'SHA-256 hash fingerprints the evidence file.',
    'IPFS upload via Pinata provides decentralized storage with CID.',
    'SQLite records off-chain metadata for efficient querying.',
    'Multiple hashes combined into Merkle tree — only root stored on-chain.',
    'Single Polygon transaction anchors entire batch via EvidenceAnchor contract.',
    'Cost reduction: 83% for 10 records, 99.7% for 100 records.',
])

# SLIDE 20: Real-Time Protection Overlay
add_title_content_slide('Real-Time Protection Overlay', [
    '2300+ line Android system overlay for continuous security monitoring.',
    'Leverages Accessibility Services to detect content in whitelisted apps.',
    'Monitors WhatsApp, Chrome, Instagram, Telegram, Facebook.',
    'Three surface modes: monitoring bubble, verdict card, call monitoring chip.',
    'URL verification against URLhaus, community blacklist, heuristic screening.',
    'Media detection captures dominant visible image for deepfake analysis.',
    'Live call monitoring with real-time voice analysis during active calls.',
])

# SLIDE 21: URL Verification & Risk Scoring
add_title_content_slide('URL Verification & Risk Scoring', [
    'Multi-layered threat intelligence for link safety assessment.',
    'Normalizes URLs, checks community blacklist, queries URLhaus API.',
    'Local heuristic screening for suspicious URL patterns.',
    'Results cached for 120 seconds with confidence scoring.',
    'Risk scoring: Call Pattern (25%), Voice (30%), Content (30%), History (15%).',
    'Risk levels: LOW (≤30), MEDIUM (31-70), HIGH (71-100).',
])

# SLIDE 22: Cybercrime Investigation Dashboard
add_title_content_slide('Cybercrime Investigation Dashboard', [
    'Web-based dashboard for law enforcement with real-time evidence management.',
    'Live evidence table with Server-Sent Events (SSE) real-time updates.',
    'Green pulsing LIVE badge showing active SSE connection status.',
    'Toast notifications with sound alerts for new evidence.',
    'Full evidence detail: SHA-256 hash, IPFS CID, Merkle proof, blockchain TX.',
    'One-click batch anchoring and direct PolygonScan links.',
    'Secure session-based authentication for investigators.',
])

# SLIDE 23: Technologies Used
slide23 = add_title_content_slide('Technologies Used', [])
tech_data = [
    ['Component', 'Technology', 'Purpose'],
    ['Mobile Frontend', 'Flutter / Dart', 'Cross-platform app with overlay'],
    ['Backend API', 'FastAPI (Python 3.11+)', 'Async REST API with Swagger UI'],
    ['AI Detection', 'NumPy, SciPy, OpenCV', 'Zero-dependency signal processing'],
    ['Cloud Models', 'HuggingFace Inference API', 'Transformer classification'],
    ['GPU Inference', 'Google Colab + ONNX', 'wav2vec2, CNN classifier'],
    ['Blockchain', 'web3.py, Polygon Amoy', 'On-chain evidence anchoring'],
    ['Smart Contract', 'Solidity', 'Merkle root storage'],
    ['IPFS', 'Pinata API', 'Decentralized evidence storage'],
    ['Dashboard', 'FastAPI + Jinja2 + SSE', 'Real-time investigator interface'],
    ['State Mgmt', 'Provider (Flutter)', 'Reactive state for protection'],
]
rows, cols = len(tech_data), len(tech_data[0])
tbl = slide23.shapes.add_table(rows, cols, Emu(457200), Emu(1800000), Emu(8229600), Emu(4800000)).table
for r in range(rows):
    for c in range(cols):
        cell = tbl.cell(r, c)
        cell.text = tech_data[r][c]
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.name = 'Times New Roman'
                run.font.size = Pt(12)
                if r == 0:
                    run.font.bold = True

# SLIDE 24: Working of the System - Image & Voice
add_title_content_slide('Working of the System', [
    'Image: Classifies input type → assigns signal weights → runs 6 signals concurrently.',
    'Failed signals excluded, weights redistributed proportionally.',
    'Weighted ensemble produces AI-generation probability with explanation.',
    'Voice: Stereo-to-mono → resample 16kHz → VAD strips silence.',
    'Real-time: 0.5s chunks processed under 200ms each.',
    'Local CPU signals combined with optional cloud wav2vec2 inference.',
])

# SLIDE 25: Working - Overlay & Evidence
add_title_content_slide('Working of the System (Cont.)', [
    'Overlay: Detects whitelisted app → activates monitoring bubble.',
    'Captures URLs for verification, media for deepfake analysis.',
    'Call monitoring: transforms into voice analysis chip during calls.',
    'Evidence: SHA-256 hash → IPFS upload → SQLite metadata.',
    'Merkle tree aggregation → single Polygon transaction anchors batch.',
    'Verification recomputes Merkle proof and compares to on-chain root.',
])

# SLIDE 26: Testing & Evaluation
add_title_content_slide('Testing and Evaluation', [
    'Unit Testing — Individual modules validated in isolation.',
    'Integration Testing — Frontend-backend-AI model communication verified.',
    'System Testing — Full application tested with real-world inputs.',
    'User Acceptance Testing — End-user evaluation of usability and UX.',
    'Performance factors: response time, processing speed, reliability assessed.',
    'Feedback-driven enhancements to UI and interaction design.',
])

# SLIDE 27: Test Cases
slide27 = add_title_content_slide('Test Case Results', [])
tc_data = [
    ['Test ID', 'Input Type', 'Expected Output', 'Result'],
    ['TC01', 'Real Image', 'Classified as Real', 'Passed'],
    ['TC02', 'Deepfake Image', 'Classified as Fake', 'Passed'],
    ['TC03', 'Real Video', 'Real with Confidence Score', 'Passed'],
    ['TC04', 'Deepfake Video', 'Fake with Confidence Score', 'Passed'],
    ['TC05', 'Invalid File', 'Error Message', 'Passed'],
]
rows, cols = len(tc_data), len(tc_data[0])
tbl = slide27.shapes.add_table(rows, cols, Emu(457200), Emu(2000000), Emu(8229600), Emu(2800000)).table
for r in range(rows):
    for c in range(cols):
        cell = tbl.cell(r, c)
        cell.text = tc_data[r][c]
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.name = 'Times New Roman'
                run.font.size = Pt(14)
                if r == 0:
                    run.font.bold = True

# SLIDE 28: Results & Performance
add_title_content_slide('Results and Performance', [
    'Image detection: 70-92% accuracy on photographs.',
    'Type-adaptive weighting prevents false positives on digital art.',
    'Voice detection: 65-85% accuracy on full audio clips.',
    'Real-time streaming support under 200ms per 0.5s chunk.',
    'Text detection: 70-88% accuracy on long-form content.',
    'Merkle batching reduces per-record costs by up to 99.7%.',
    'Local image analysis completes under 150ms.',
])

# SLIDE 29: Conclusion
add_title_content_slide('Conclusion', [
    'RiskGuard integrates multi-modal detection, blockchain evidence, and real-time protection.',
    'Production-grade deepfake detection achievable on commodity hardware.',
    '6-signal image ensemble with type-adaptive weighting prevents false positives.',
    'Voice pipeline achieves real-time streaming under 200ms per chunk.',
    'Blockchain evidence system reduces costs by 99.7% through Merkle batching.',
    'Real-time overlay provides continuous background monitoring without user intervention.',
    'Dual-interface successfully serves both citizens and law enforcement.',
])

# SLIDE 30: Future Scope
add_title_content_slide('Future Scope', [
    'Deploy fine-tuned quantized models (wav2vec2 + ViT) locally via ONNX Runtime.',
    'Adversarial robustness testing against perturbation-based evasion attacks.',
    'Cross-modal fusion combining image, text, and voice for multi-modal documents.',
    'Migration from Polygon Amoy testnet to Polygon mainnet or zkEVM.',
    'Extend overlay system to support iOS platform.',
    'Enhance SMS phishing detection with on-device NLP models.',
    'Continuous model retraining as new generative AI techniques emerge.',
])

# SLIDE 31: References
add_title_content_slide('References', [
    '1. Cozzolino & Verdoliva (2020) — Noiseprint: CNN-Based Camera Model Fingerprint.',
    '2. Tak et al. (2021) — End-to-End Anti-Spoofing with RawNet2. ICASSP.',
    '3. Baevski et al. (2020) — wav2vec 2.0: Self-Supervised Learning. NeurIPS.',
    '4. Hans et al. (2024) — Spotting LLMs with Binoculars. ICML.',
    '5. Sahidullah et al. (2015) — Block Level Features for Anti-Spoofing.',
    '6. ASVspoof Consortium (2024) — ASVspoof 5. ISCA Interspeech.',
    '7. Benet, J. (2014) — IPFS: Content Addressed P2P File System.',
])

# SLIDE 32: Thank You
slide_ty = prs.slides.add_slide(content_layout)
# Clear content placeholder
if 1 in slide_ty.placeholders:
    sp = slide_ty.placeholders[1]._element
    sp.getparent().remove(sp)
tf = slide_ty.placeholders[0].text_frame
tf.clear()
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = 'THANK YOU'
run.font.name = 'Times New Roman'
run.font.size = Pt(60)

prs.save(OUTPUT)
print(f"✅ RiskGuard Final Review PPT saved to: {OUTPUT}")
print(f"   Total slides: {len(prs.slides)}")
